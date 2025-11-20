#!/usr/bin/env python3
"""
Generate PhD Research Synopsis Presentation
Title: From Prediction to Performance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Create the PhD synopsis presentation"""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    DARK_BLUE = RGBColor(25, 55, 95)
    ACCENT_BLUE = RGBColor(0, 112, 192)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    DARK_GRAY = RGBColor(89, 89, 89)
    WHITE = RGBColor(255, 255, 255)
    ORANGE = RGBColor(237, 125, 49)

    def add_title_slide(title, subtitle):
        """Add title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = DARK_BLUE
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

        return slide

    def add_section_slide(title):
        """Add section divider slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(3), prs.slide_width, Inches(1.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_BLUE
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.15), Inches(11.333), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, bullets, notes=""):
        """Add standard content slide with bullets"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = DARK_BLUE
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.833), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Handle sub-bullets
            if bullet.startswith("  - "):
                p.text = bullet[4:]
                p.level = 1
                p.font.size = Pt(18)
            else:
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)

            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
        """Add two-column comparison slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = DARK_BLUE
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Left column title
        left_title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.5), Inches(0.6))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        # Left column content
        left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.5), Inches(4.8))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(left_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)

        # Right column title
        right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.6))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        # Right column content
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.8))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(right_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)

        return slide

    def add_highlight_slide(title, main_text, subtitle=""):
        """Add highlight/quote slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = LIGHT_GRAY
        background.line.fill.background()

        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(0.1), Inches(2.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_BLUE
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Main text
        main_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.833), Inches(2.5))
        tf = main_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main_text
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.833), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = ACCENT_BLUE

        return slide

    # ========== CREATE SLIDES ==========

    # Slide 1: Title
    add_title_slide(
        "From Prediction to Performance",
        "Structuring Founding Team Decisions in AI-Supported Entrepreneurship\n\nMr. Syed Faraaz | PhD Synopsis | Nirwan University Jaipur"
    )

    # Slide 2: The Problem (Hook)
    add_highlight_slide(
        "The Problem",
        "I've watched founding teams fall apart over a number on a screen.",
        "A score of 74 vs 76 can determine funding access—and team dynamics"
    )

    # Slide 3: What I Observed
    add_content_slide(
        "What Actually Happens",
        [
            "Team scores 74 (one point below threshold):",
            "  - Everyone starts questioning the inputs",
            "  - 'Should we have said three advisors instead of two?'",
            "  - Tension, blame, gaming behavior",
            "",
            "Team scores 76 (barely above threshold):",
            "  - Treat it as gospel truth",
            "  - Ignore warning signs in burn rate",
            "  - Overconfidence and groupthink",
            "",
            "The technology works. It's what happens AFTER that's messy."
        ]
    )

    # Slide 4: Section - Research Context
    add_section_slide("Research Context: Flash CAMP")

    # Slide 5: Flash System Overview
    add_content_slide(
        "The Flash CAMP System",
        [
            "Production ML platform for startup assessment:",
            "  - 22 machine learning models (stage & industry-specific)",
            "  - 126 engineered features from 45 inputs",
            "  - 81% prediction accuracy",
            "  - 6,500+ assessed startups with tracked outcomes",
            "",
            "CAMP Framework:",
            "  - Capital: Runway, burn, efficiency",
            "  - Advantage: Moats, IP, differentiation",
            "  - Market: TAM/SAM/SOM, timing, growth",
            "  - People: Team experience, gaps, advisors"
        ]
    )

    # Slide 6: The Multi-Phase Journey
    add_content_slide(
        "The Multi-Phase Decision Journey",
        [
            "Phase 1: Assessment",
            "  - Founder fills detailed forms (30-45 min)",
            "",
            "Phase 2: Results Reveal",
            "  - Overall score → 75+ unlocks funding access",
            "  - CAMP dimensional breakdown",
            "",
            "Phase 3: Strategic Report",
            "  - 50,000-word McKinsey-level analysis",
            "",
            "Phase 4: AI Advisor",
            "  - Conversational AI trained on 10,000+ real cases",
            "",
            "Each phase = new decision moment. How teams handle this matters."
        ]
    )

    # Slide 7: The 75-Point Threshold
    add_highlight_slide(
        "The Natural Experiment",
        "Score ≥ 75 = Access to funding networks\nScore < 75 = No access",
        "This binary gate creates real stakes and fascinating behavioral dynamics"
    )

    # Slide 8: Section - The Gap
    add_section_slide("What We Don't Know")

    # Slide 9: Literature Gaps
    add_two_column_slide(
        "Literature vs. Reality",
        "What Research Studies",
        [
            "Individual decision-makers",
            "Standardized operational tasks",
            "Single algorithm, single decision",
            "Lab settings with students",
            "Hypothetical scenarios"
        ],
        "What Actually Happens",
        [
            "Founding teams (2-5 people)",
            "Ambiguous strategic decisions",
            "Multi-phase AI engagement",
            "Real stakes (survival, funding)",
            "History, power dynamics, emotions"
        ]
    )

    # Slide 10: Key Research Gaps
    add_content_slide(
        "Critical Unknowns",
        [
            "How do teams actually integrate algorithmic forecasts?",
            "",
            "Does structured process help or slow things down?",
            "",
            "How do AI 'verdicts' affect team dynamics?",
            "  - Who speaks up? Who stays quiet?",
            "",
            "What happens near threshold points (73-77)?",
            "",
            "Does any of this predict venture performance?"
        ]
    )

    # Slide 11: Section - Research Design
    add_section_slide("Research Questions & Design")

    # Slide 12: Core Questions
    add_content_slide(
        "Four Driving Questions",
        [
            "RQ1: Process → How do teams currently structure decisions with AI forecasts?",
            "",
            "RQ2: Dynamics → Do structured processes affect psychological safety and voice?",
            "",
            "RQ3: Threshold → What happens behaviorally near the 75-point cutoff?",
            "",
            "RQ4: Outcomes → Can we link process quality to venture performance?"
        ]
    )

    # Slide 13: Hypotheses Overview
    add_content_slide(
        "Key Hypotheses",
        [
            "Structure & Integration",
            "  - H1: Structured teams show balanced AI integration (not blind trust or rejection)",
            "",
            "Team Dynamics",
            "  - H3: Structure increases psychological safety",
            "  - H4: Structure increases voice behavior",
            "",
            "Threshold Effects",
            "  - H5: Teams near 75 behave differently",
            "  - H6: Structure matters MORE when stakes are higher",
            "",
            "Performance",
            "  - H8: Structured process → better venture outcomes"
        ]
    )

    # Slide 14: Section - Methodology
    add_section_slide("Methodology: Three Phases")

    # Slide 15: Phase I
    add_content_slide(
        "Phase I: Qualitative Foundation (Months 1-9)",
        [
            "Purpose: Understand territory before measuring it",
            "",
            "Sample: 10-12 ventures using Flash CAMP",
            "  - Various scores, including 73-77 threshold cases",
            "",
            "Data Collection:",
            "  - Semi-structured interviews (3-4 per venture)",
            "  - Observation of actual decision meetings",
            "  - Document analysis (inputs, outputs, memos)",
            "",
            "Output: Process typology, preliminary measurement instrument"
        ]
    )

    # Slide 16: Phase II
    add_content_slide(
        "Phase II: Experimental Test (Months 13-21)",
        [
            "Purpose: Establish causal effects",
            "",
            "Design: 2×2 Factorial",
            "  - Process: Structured vs Unstructured",
            "  - Threshold: Near (74/76) vs Far (68/82)",
            "",
            "Sample: 80 teams (240-320 participants)",
            "",
            "Structured Protocol Includes:",
            "  - Independent pre-judgments before reveal",
            "  - Divided report analysis (each person = one CAMP dimension)",
            "  - Required dissent round",
            "  - AI query logging with collective review"
        ]
    )

    # Slide 17: Phase III
    add_content_slide(
        "Phase III: Field Study (Months 19-30)",
        [
            "Purpose: Test at scale with real outcomes",
            "",
            "Sample: 250-300 teams from Flash database",
            "",
            "Data Sources:",
            "  - Survey: Decision process, psych safety, voice",
            "  - Archival: 126 features, model predictions, usage",
            "  - Outcomes: Survival, funding, pivots, score changes",
            "",
            "Analysis:",
            "  - Regression, survival analysis, panel models",
            "  - Threshold regression discontinuity (74 vs 76)",
            "  - Propensity score matching for selection bias"
        ]
    )

    # Slide 18: Why This Design
    add_two_column_slide(
        "Methodological Strength",
        "Triangulation",
        [
            "Phase I: Rich qualitative insight",
            "Phase II: Causal claims",
            "Phase III: Scale + real outcomes",
            "",
            "Each phase builds on previous",
            "Can reconcile contradictions"
        ],
        "Unique Advantages",
        [
            "Real stakes, real teams",
            "Actual outcome tracking",
            "Natural experiment at 75",
            "Multi-phase AI engagement",
            "Production system, not prototype"
        ]
    )

    # Slide 19: Section - Contributions
    add_section_slide("Expected Contributions")

    # Slide 20: Theoretical
    add_content_slide(
        "Theoretical Contributions",
        [
            "Behavioral Strategy",
            "  - Team-level process governance (not just individual debiasing)",
            "",
            "Organizational Behavior",
            "  - How AI verdicts affect psychological safety and voice",
            "",
            "Entrepreneurship",
            "  - Evidence-based protocols for AI-supported strategy",
            "",
            "Human-AI Interaction",
            "  - Sociotechnical governance frameworks for multi-phase AI"
        ]
    )

    # Slide 21: Practical
    add_content_slide(
        "Practical Contributions",
        [
            "For Founding Teams:",
            "  - Usable protocols, checklists, templates",
            "",
            "For AI Platforms (Flash):",
            "  - Pre-commitment interface",
            "  - Structured report navigator",
            "  - Conversation logging + team review",
            "  - Decision documentation templates",
            "",
            "For Investors & Accelerators:",
            "  - Due diligence on decision processes",
            "  - Training curricula"
        ]
    )

    # Slide 22: Timeline
    add_content_slide(
        "Timeline (36 Months)",
        [
            "Year 1: Foundation",
            "  - Literature, ethics, Phase I qualitative work",
            "  - Instrument development",
            "",
            "Year 2: Testing",
            "  - Phase II experiments (80 teams)",
            "  - Phase III survey deployment",
            "",
            "Year 3: Analysis & Writing",
            "  - Outcome data, statistical analysis",
            "  - Tool development, dissertation",
            "",
            "Publications: 3-4 papers (qual, experimental, field study)"
        ]
    )

    # Slide 23: Why This Matters
    add_highlight_slide(
        "Why This Matters",
        "AI tools for strategic decisions are proliferating. The question is shifting from 'what does the AI say?' to 'how do we decide what to do about what the AI says?'",
        "This research provides evidence-based answers"
    )

    # Slide 24: The Ask
    add_content_slide(
        "Summary",
        [
            "Research Question:",
            "  - How should founding teams structure decisions with AI forecasts?",
            "",
            "Unique Opportunity:",
            "  - Direct access to Flash CAMP (6,500+ ventures, real outcomes)",
            "  - Natural experiment at 75-point threshold",
            "",
            "Contributions:",
            "  - Theory: Extend behavioral strategy to AI-augmented teams",
            "  - Practice: Usable tools embedded in production platform",
            "",
            "Outcome: Evidence-based governance for AI-augmented entrepreneurship"
        ]
    )

    # Slide 25: Closing
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Thank you
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Questions
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Mr. Syed Faraaz | Nirwan University Jaipur"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = "/home/user/nuj-synopsis-presentation/synopsis_presentation.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
